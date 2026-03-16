#!/usr/bin/env python3
"""
Export aurora-expenses docs to PDF and PPTX.

- docs/aurora-school-strategy-2026.html  → docs/aurora-school-strategy-2026.pdf
- docs/aurora-growth-strategy-deck.html  → docs/aurora-growth-strategy-deck.pptx
"""
import os, time
from pathlib import Path
from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from PIL import Image
import io

DOCS = Path(__file__).parent.parent / "docs"
ANALYSIS_HTML = DOCS / "aurora-school-strategy-2026.html"
DECK_HTML     = DOCS / "aurora-growth-strategy-deck.html"
ANALYSIS_PDF  = DOCS / "aurora-school-strategy-2026.pdf"
DECK_PPTX     = DOCS / "aurora-growth-strategy-deck.pptx"
TOTAL_SLIDES  = 15
SLIDE_W_PX    = 1920
SLIDE_H_PX    = 1080


def export_analysis_pdf(page):
    print("→ Exporting analysis to PDF…")
    page.goto(f"file://{ANALYSIS_HTML.resolve()}", wait_until="networkidle")
    page.wait_for_timeout(1500)
    page.pdf(
        path=str(ANALYSIS_PDF),
        format="A4",
        margin={"top": "18mm", "bottom": "18mm", "left": "16mm", "right": "16mm"},
        print_background=True,
        prefer_css_page_size=False,
    )
    size = ANALYSIS_PDF.stat().st_size // 1024
    print(f"   ✓ {ANALYSIS_PDF.name}  ({size} KB)")


def screenshot_slide(page, index: int) -> bytes:
    """Navigate to a slide by index and return a PNG screenshot."""
    page.evaluate(f"goTo({index})")
    page.wait_for_timeout(900)   # let reveal animations settle
    return page.screenshot(full_page=False)


def build_pptx(screenshots: list[bytes]):
    print("→ Building PPTX…")
    prs = Presentation()

    # 16:9 widescreen — 10 in × 5.625 in
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(5.625)

    blank_layout = prs.slide_layouts[6]   # completely blank layout

    for i, png_bytes in enumerate(screenshots):
        slide = prs.slides.add_slide(blank_layout)

        # Fill slide background black (fallback if image has gaps)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0x1A, 0x17, 0x14)

        # Add screenshot as full-slide image
        img_stream = io.BytesIO(png_bytes)
        pic = slide.shapes.add_picture(
            img_stream,
            left=Emu(0),
            top=Emu(0),
            width=prs.slide_width,
            height=prs.slide_height,
        )
        # Send image to back
        slide.shapes._spTree.remove(pic._element)
        slide.shapes._spTree.insert(2, pic._element)

        print(f"   slide {i+1:02d}/{TOTAL_SLIDES} added")

    prs.save(str(DECK_PPTX))
    size = DECK_PPTX.stat().st_size // 1024
    print(f"   ✓ {DECK_PPTX.name}  ({size} KB)")


def export_deck_pptx(page):
    print("→ Loading deck and taking slide screenshots…")
    page.set_viewport_size({"width": SLIDE_W_PX, "height": SLIDE_H_PX})
    page.goto(f"file://{DECK_HTML.resolve()}", wait_until="networkidle")
    # Let fonts and animations initialise
    page.wait_for_timeout(2000)

    screenshots = []
    for i in range(TOTAL_SLIDES):
        png = screenshot_slide(page, i)
        screenshots.append(png)
        print(f"   screenshot {i+1:02d}/{TOTAL_SLIDES}")

    build_pptx(screenshots)


def main():
    DOCS.mkdir(exist_ok=True)
    with sync_playwright() as pw:
        browser = pw.chromium.launch(headless=True)

        # --- PDF (analysis) ---
        ctx_pdf = browser.new_context()
        page_pdf = ctx_pdf.new_page()
        export_analysis_pdf(page_pdf)
        ctx_pdf.close()

        # --- PPTX (deck) ---
        ctx_deck = browser.new_context(
            viewport={"width": SLIDE_W_PX, "height": SLIDE_H_PX},
        )
        page_deck = ctx_deck.new_page()
        export_deck_pptx(page_deck)
        ctx_deck.close()

        browser.close()

    print("\n✅ Done.")
    print(f"   PDF  → {ANALYSIS_PDF}")
    print(f"   PPTX → {DECK_PPTX}")


if __name__ == "__main__":
    main()
